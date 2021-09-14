from typing import List
from pptx import Presentation

import requests
import uuid
import logging
import asyncio
import threading, queue
import os
import json
import pika
import time
import shutil

categ = {'W-18+': 'F18', 
        'W-25+': 'F25',
        'W-30+': 'F30',
        'W-35+': 'F35',
        'W-40+': 'F40',
        'W-45+': 'F45',
        'W-50+': 'F50',
        'W-55+': 'F55',
        'W-65+': 'F65',
        'W-70+': 'F70',
        'W-75+': 'F75',
        'M-18+': 'M18', 
        'M-25+': 'M25',
        'M-30+': 'M30',
        'M-35+': 'M35',
        'M-40+': 'M40',
        'M-45+': 'M45',
        'M-50+': 'M50',
        'M-55+': 'M55',
        'M-65+': 'M65',
        'M-70+': 'M70',
        'M-75+': 'M75'  
        }

print(categ['M-75+'])

def search_and_replace(search_str, repl_str, input, output):
    #prs = Presentation(input)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(search_str))!=-1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str))
                    text_frame.paragraphs[0].runs[0].text = new_text
    #prs.save(output)


def search_user(category, rank, inData):
    for course in inData['Course_101']:
        #print(course)
        if(course['category'] == category and course['FINISH_ET1_RK_category'] == rank):
            return course['first'],course['last'],course['FINISH_ET1_Time'],course['FINISH_ET1_RK_category']
        else:
            return 'NON','NON','NON','NON' #what to do in that case maybe name and no time?

def replace_key(key, index, type):
        KEY = categ[key]+"_"+str(index)+type
        print(KEY)
        first, last, ttime, category = search_user(key,str(ind),data)
        print( first+' '+last,ttime,category)
        if(type == '_TIME'):
            search_and_replace(KEY,ttime,'OUT.pptx','OUT.pptx' ) 
        else:
            search_and_replace(KEY,first+' '+last,'OUT.pptx','OUT.pptx' ) 

print("start...")
shutil.copy2('_Presentation_AG-1-3.pptx', 'OUT.pptx') 

prs = Presentation('OUT.pptx')

r = requests.get("http://win2.fh-timing.com:8081/result/json?course=101&detail=start,first,last,category,club&splitnr=199&showaw=2")
data = r.json()

for key in categ:
    for ind in range(1, 4):
        replace_key(key, ind, '_NAME') # FIRST Replace Name
        replace_key(key, ind, '_TIME') # SECOND Replace Rime

prs.save('OUT.pptx') 
#search_and_replace('F40_1_TIME','AAbbbA','OUT.pptx','OUT.pptx' )
print("... done!")
