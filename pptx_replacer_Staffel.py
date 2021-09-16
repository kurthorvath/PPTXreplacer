#from typing import List
from pptx import Presentation

import requests
#import uuid
#import logging
#import asyncio
#import threading, queue
#import os
#import json
#import pika
#import time
import shutil

categ = {'RX': 'X', 
        'RF': 'F',
        'RM': 'M'
        }

def search_and_replace(search_str, team, teamtime, sName, sTime, bName, bTime, rName, rTime,  input, output):
    prs = Presentation(input)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                #print("each",shape.text, "search4: ", search_str)
                if(shape.text.find(search_str+'_FULL_NAME'))!=-1:
                    text_frame = shape.text_frame
                    print("each -> FOUND:",shape.text, "search4: ", search_str+'_FULL_NAME', "->", team)
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str+'_FULL_NAME'), str(team))
                    text_frame.paragraphs[0].runs[0].text = new_text
                
                    cur_text = text_frame.paragraphs[1].runs[0].text
                    new_text = cur_text.replace(str(search_str+'_FULL_TIME'), str(teamtime))
                    text_frame.paragraphs[1].runs[0].text = new_text

                    cur_text = text_frame.paragraphs[2].runs[0].text
                    new_text = cur_text.replace(str(search_str+'_SWIM_NAME'), str(sName))
                    text_frame.paragraphs[2].runs[0].text = new_text

                    cur_text = text_frame.paragraphs[2].runs[0].text
                    new_text = cur_text.replace(str(search_str+'_SWIM_TIME'), str(sTime))
                    text_frame.paragraphs[2].runs[0].text = new_text

                    cur_text = text_frame.paragraphs[2].runs[1].text
                    new_text = cur_text.replace(str(search_str+'_BIKE_NAME'), str(bName))
                    text_frame.paragraphs[2].runs[1].text = new_text

                    cur_text = text_frame.paragraphs[2].runs[1].text
                    new_text = cur_text.replace(str(search_str+'_BIKE_TIME'), str(bTime))
                    text_frame.paragraphs[2].runs[1].text = new_text

                    cur_text = text_frame.paragraphs[2].runs[2].text
                    new_text = cur_text.replace(str(search_str+'_RUN_NAME'), str(rName))
                    text_frame.paragraphs[2].runs[2].text = new_text

                    cur_text = text_frame.paragraphs[2].runs[2].text
                    new_text = cur_text.replace(str(search_str+'_RUN_TIME'), str(rTime))
                    text_frame.paragraphs[2].runs[2].text = new_text


    prs.save(output)


def search_user(category, rank, inData):
    print(len(inData['Course_1']))
    for index in range(0, len(data['Course_1'])):
        if(data['Course_1'][index]['category'] == category and data['Course_1'][index]['Finish_RK_Agegroup'] == rank):
            print("FOUND",data['Course_1'][index])
            return data['Course_1'][index]['last'],data['Course_1'][index]['Finish_Time'],data['Course_1'][index]['rem2'],data['Course_1'][index]['Swim_Time'],data['Course_1'][index]['rem4'],data['Course_1'][index]['Bike Leg_Time'],data['Course_1'][index]['rem6'],data['Course_1'][index]['Run Leg_Time']        
    return 'NON','NON','NON','NON','NON','NON','NON','NON' #what to do in that case maybe name and no time?
    
def replace_key(key, index):
        #KEY = categ[key]+"_"+str(index)
        #print(key, '->', KEY, len(data['Course_1']))

        team, teamtime, sName, sTime, bName, bTime, rName, rTime = search_user(key,str(index),data)
        print( key, index , ">",team, teamtime, sName, sTime, bName, bTime, rName, rTime)
        
        KEY = 'REL_'+categ[key]+"_"+str(index)
        print("KEY -><",KEY)
        search_and_replace(KEY,team, teamtime, sName, sTime, bName, bTime, rName, rTime ,'OUT_RELAY.pptx','OUT_RELAY.pptx' ) 





print("start...")
shutil.copy2('_Presentation_Relay-1-3.pptx', 'OUT_RELAY.pptx') 
#prs = Presentation('OUT_RELAY.pptx')
r = requests.get("http://192.168.70.52:8081/result/json?course=1&detail=start,first,last,category,rem2,rem4,rem6&splitnr=599110,199110,399299,599499&showaw=2")
data = r.json()
print(len(data['Course_1']))


for key in categ:
    for ind in range(1, 4):
        replace_key(key, ind) # FIRST Replace Name

#prs.save('OUT_RELAY.pptx') 
## for testing
#search_and_replace('REL_X_1_FULL_NAME','ABCDEFG', '', '', '', '', '', '', '' ,'OUT_RELAY.pptx','OUT_RELAY.pptx' ) 
#search_and_replace('REL_X_2_FULL_NAME','ABCDEFG', '', '', '', '', '', '', '' ,'OUT_RELAY.pptx','OUT_RELAY.pptx' ) 
#search_and_replace('REL_X_3_FULL_NAME','ABCDEFG', '', '', '', '', '', '', '' ,'OUT_RELAY.pptx','OUT_RELAY.pptx' ) 
print("... done!")
