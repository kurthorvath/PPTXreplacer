# -*- coding: utf-8 -*-
"""
Created on Wed Sep 15 13:57:59 2021

@author: uran
"""

from pptx import Presentation
import requests


TEMPLATE = '..\\presentations\\_Presentation_AG-1-3.pptx'
OUT = '..\\presentations\\_Presentation_AG-1-3_ChU.pptx'
CATEGORIES = {'F18': 'F18-24', 
        'F25': 'F25-29',
        'F30': 'F30-34',
        'F35': 'F35-39',
        'F40': 'F40-44',
        'F45': 'F45-49',
        'F50': 'F50-54',
        'F55': 'F55-59',
        'F60': 'F60-64',
        'F65': 'F65-69',
        'F70': 'F70-74',
        'F75': 'F75-79',
        'M18': 'M18-24', 
        'M25': 'M25-29',
        'M30': 'M30-34',
        'M35': 'M35-39',
        'M40': 'M40-44',
        'M45': 'M45-49',
        'M50': 'M50-54',
        'M55': 'M55-59',
        'M60': 'M60-64',
        'M65': 'M65-69',
        'M70': 'M70-74',
        'M75': 'M75-79'
        }
COURSE = 1

r = requests.get("http://192.168.70.52:8081/result/json?course=1&detail=start,first,last,category,club&splitnr=599110&showaw=2")
data = r.json()
data = data['Course_%s' % COURSE]

presentation = Presentation(TEMPLATE)

for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for run in shape.text_frame.paragraphs[0].runs:
                for cat in CATEGORIES.keys():
                    if cat in run.text:
                        for rk in [r+1 for r in range(3)]:
                            if ('%s_%s' % (cat, rk)) in run.text:
                                for kind in ['NAME', 'TIME']:
                                    if ('%s_%s_%s' % (cat, rk, kind)) in run.text:
                                        run.text = run.text.replace(str('%s_%s_%s' % (cat, rk, kind)), str('found you'))

presentation.save(OUT)
